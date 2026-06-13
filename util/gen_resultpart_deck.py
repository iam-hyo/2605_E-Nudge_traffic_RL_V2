"""gen_resultpart_deck.py — E-Nudge 결과파트 덱 (design_guide_v4 적용).

E-Nudge_resultPart_guide.pdf 구조 재제작 + 빈 '경로3' 자리 채움 + 12x12 일반화 슬라이드 추가.
데이터: output/08_1412_gangnam_v2 (강남 v2) + 09_0119 (12x12 held-out 다중학습).

design_guide_v4:
  16:9 13.333"x7.5" · light mode · 1 primary accent(#2D6CDF) · 2nd accent(#E23B3B, 결론/경고만)
  헤더 앵커(타이틀 마커|+제목 / 메시지바 / 메타행) · 카드 hairline · 폰트 1종(Pretendard)
사용: venv/bin/python util/gen_resultpart_deck.py
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
FIG  = ROOT / "output" / "08_1412_gangnam_v2"
OUT  = FIG / "E-Nudge_ResultPart_v4.pptx"

# ── v4 토큰 ──
CANVAS  = RGBColor(0xFF,0xFF,0xFF)
PARCH   = RGBColor(0xF5,0xF5,0xF7)
CARD    = RGBColor(0xF5,0xF8,0xFC)
DARK    = RGBColor(0x1B,0x2A,0x4A)
ACCFILL = RGBColor(0xEA,0xF1,0xFC)
ACCENT  = RGBColor(0x2D,0x6C,0xDF)
ACC2    = RGBColor(0xE2,0x3B,0x3B)
INK     = RGBColor(0x1F,0x29,0x37)
ONDARK  = RGBColor(0xFF,0xFF,0xFF)
MUTED   = RGBColor(0x5B,0x64,0x72)
FAINT   = RGBColor(0x9C,0xA3,0xAF)
LINE    = RGBColor(0xE2,0xE7,0xEF)
GREEN   = RGBColor(0x16,0xA8,0x58)
FONT = "Pretendard"
CW, CH = 13.333, 7.5
EMU_IN = 914400

def _f(run, size, color, bold=False):
    run.font.name=FONT; run.font.size=Pt(size); run.font.bold=bold
    run.font.color.rgb=color
    rPr=run._r.get_or_add_rPr()
    rPr.append(rPr.makeelement(qn('a:ea'),{'typeface':FONT}))

def rect(s,x,y,w,h,fill,line=None,lw=1.0,rounded=False,radius=0.08):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                           Inches(x),Inches(y),Inches(w),Inches(h))
    if rounded:
        try: shp.adjustments[0]=radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False
    return shp

def text(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sa=2,ls=1.0,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap; tf.vertical_anchor=anchor
    if runs and not isinstance(runs[0],list): runs=[runs]
    for i,pr in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.space_before=Pt(0); p.line_spacing=ls
        for (t,sz,c,b) in pr:
            r=p.add_run(); r.text=t; _f(r,sz,c,b)
    return tb

def slide(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6]); rect(s,0,0,CW,CH,CANVAS); return s

def header(s, kicker, title, msg=None, meta=None):
    # 타이틀 마커 | + 제목
    rect(s,0.31,0.30,0.06,0.42,ACCENT)
    text(s,0.30,0.16,11,0.3,[(kicker,11,ACCENT,True)])
    text(s,0.46,0.40,12.4,0.62,[(title,25,INK,True)])
    if msg: text(s,0.46,1.02,12.3,0.34,msg if isinstance(msg,list) else [(msg,13.5,MUTED,False)])
    if meta: text(s,0.46,1.36,12.3,0.26,[(meta,11,FAINT,False)])

def footer(s,n):
    text(s,0.4,7.06,9,0.26,[("Fuel-Optimal Routing · 강화학습 기반 연료소비 최소화 경로 안내",9,FAINT,False)])
    text(s,11.9,7.06,1.0,0.26,[(f"{n} / 9",9,FAINT,False)],align=PP_ALIGN.RIGHT)

def badge(s,x,y,num,color=ACCENT,d=0.34):
    b=rect(s,x,y,d,d,color,rounded=True,radius=0.5)
    tf=b.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=str(num); _f(r,12,ONDARK,True); return b

def pic(s,path,x,y,w=None,h=None):
    if Path(path).exists():
        kw={}
        if w: kw["width"]=Inches(w)
        if h: kw["height"]=Inches(h)
        s.shapes.add_picture(str(path),Inches(x),Inches(y),**kw)

def simple_table(s, x, y, w, col_w, rows, fontsz=10.5, header_dark=True, rh=0.34):
    """rows[0]=header. col_w 비율 리스트. 셀=(text,color,bold) 또는 str."""
    cx=x
    cols=len(col_w); tot=sum(col_w)
    yy=y
    for ri,row in enumerate(rows):
        cx=x
        for ci,cell in enumerate(row):
            cwi=w*col_w[ci]/tot
            if ri==0:
                fill=DARK if header_dark else ACCFILL
            else:
                fill=CARD if ri%2 else CANVAS
            rect(s,cx,yy,cwi,rh,fill,line=LINE,lw=0.75)
            default_c=(ONDARK if ri==0 and header_dark else INK)
            if isinstance(cell,tuple):
                t,c,b=cell
                if not isinstance(c,RGBColor): c=default_c
            else: t,c,b=cell,default_c,(ri==0)
            al=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
            text(s,cx+0.06,yy,cwi-0.1,rh,[(t,fontsz,c,b)],align=al,anchor=MSO_ANCHOR.MIDDLE)
            cx+=cwi
        yy+=rh
    return yy

# ════════════════════════════════════════════════════════════════
def build():
    prs=Presentation(); prs.slide_width=Emu(int(CW*EMU_IN)); prs.slide_height=Emu(int(CH*EMU_IN))

    # ===== 1. 실험 환경 구축 =====
    s=slide(prs)
    header(s,"07 · 실험 환경 구축","강남구 도로망 기반 시뮬레이션 환경 구축",
           [("합성 환경에서 더 나아가 DQN 방법론의 ",13.5,MUTED,False),("실제 도로 일반화 가능성",13.5,INK,True),("을 탐색한다.",13.5,MUTED,False)])
    # 좌: 지도 카드
    rect(s,0.4,1.7,5.0,5.05,CARD,line=LINE,rounded=True)
    text(s,0.6,1.85,4.6,0.3,[("그림 1. 강남구 도로망 구조",12,INK,True)])
    text(s,0.6,2.15,4.6,0.25,[("노드=교차로 · 링크=도로구간",9.5,MUTED,False)])
    pic(s,FIG/"gn_od1_sinnonhyeon_suseo_fullmap.png",0.7,2.45,w=4.4)
    rect(s,0.62,2.5,1.45,0.95,CANVAS,line=LINE,rounded=True)
    text(s,0.72,2.56,1.3,0.9,[[("노드 1995",10,INK,True)],[("링크 2439",10,INK,True)],[("신호 231",10,INK,True)]],ls=1.15)
    # 우: 3단계 + 가정
    text(s,5.7,1.78,7.2,0.3,[("실제 도로망을 노드·링크 환경으로 변환",13.5,ACCENT,True)])
    steps=[("1","도로 데이터 추출","QGIS로 강남구 노드·링크 CSV 추출"),
           ("2","노드·링크 구조화","교차로=노드, 도로구간=링크"),
           ("3","시뮬레이션 환경","노드·링크를 JSON 형태로 변환")]
    yy=2.25
    for n,t,d in steps:
        badge(s,5.7,yy,n); text(s,6.16,yy-0.02,6.6,0.3,[(t,13,INK,True)])
        text(s,6.16,yy+0.30,6.6,0.3,[(d,11,MUTED,False)]); yy+=0.86
    rect(s,5.7,4.95,7.2,1.75,DARK,rounded=True)
    text(s,6.0,5.12,6.7,0.34,[("시뮬레이션 실행 가정",12.5,RGBColor(0x8F,0xB4,0xF5),True)])
    g=[("1","현재 노드","위치·시간·신호·인접링크 속도"),("2","다음 노드 선택","연결된 인접 노드 중 1개"),("3","주행 결과 계산","이동·대기시간·연료 누적")]
    gx=6.0
    for n,t,d in g:
        rect(s,gx,5.5,2.18,1.0,CANVAS,line=None,rounded=True)
        text(s,gx+0.14,5.6,1.9,0.3,[(f"{n}  ",11,ACCENT,True),(t,11,INK,True)])
        text(s,gx+0.14,5.9,1.95,0.5,[(d,9.5,MUTED,False)],ls=1.05); gx+=2.34
    footer(s,1)

    # ===== 2. 강남 실데이터 환경 및 경로 =====
    s=slide(prs)
    header(s,"07 · 실험 환경 구축","강남 실데이터 환경 및 경로",
           [("출근시간 통행 동선 ",13.5,MUTED,False),("OD 3쌍",13.5,INK,True),("과 ",13.5,MUTED,False),("학습 파라미터",13.5,INK,True),("를 소개한다.",13.5,MUTED,False)])
    # 좌: OD-1,2 표
    text(s,0.4,1.72,5.2,0.3,[("경로 환경 (출근 간선축)",13,ACCENT,True)])
    simple_table(s,0.4,2.1,5.2,[1.0,1.5],[
        [("OD-1","",True),("신논현 → 수서","",True)],
        [("직선/주행","",False),("5.6km / 9.2km","",False)],
        [("성격","",False),("강남대로 정체 vs 우회","",False)]],fontsz=10.5,rh=0.42)
    simple_table(s,0.4,3.75,5.2,[1.0,1.5],[
        [("OD-2","",True),("양재 → 영동대교","",True)],
        [("직선/주행","",False),("4.7km / 6.4km","",False)],
        [("성격","",False),("강남대로·영동대로 북상","",False)]],fontsz=10.5,rh=0.42)
    simple_table(s,0.4,5.4,5.2,[1.0,1.5],[
        [("OD-3","",True),("세곡 → 삼성(GBD)","",True)],
        [("직선/주행","",False),("5.9km / 7.6km","",False)],
        [("성격","",False),("주거→업무지구, 신호밀집","",False)]],fontsz=10.5,rh=0.42)
    # 우: 경로 이미지 + 파라미터
    pic(s,FIG/"gn_od1_sinnonhyeon_suseo_fullmap.png",6.0,1.7,h=2.7)
    pic(s,FIG/"gn_od2_yangjae_yeongdong_fullmap.png",9.7,1.7,h=2.7)
    text(s,6.0,4.32,7,0.25,[("OD-1 신논현→수서                          OD-2 양재→영동대교  (파랑=거리최단, 주황=연료최단)",9,MUTED,False)])
    text(s,6.0,4.62,7,0.3,[("실험 환경 vs 실제 환경 — 파라미터 조정",13,ACCENT,True)])
    simple_table(s,6.0,4.98,6.9,[1.5,0.8,0.8,2.0],[
        ["항목","12x12","강남","조정 이유"],
        [("shaping_weight","",False),"500","1500",("step↑·곡률 보정","",False)],
        [("arrival_bonus","",False),"200","700",("연료 1500~2500mL 비율","",False)],
        [("train_max_steps","",False),"220","400",("평균 경로 길이 비례","",False)],
        [("episodes","",False),("4~6k","",False),("15k","",False),("큰 환경 학습 길게","",False)]],fontsz=9.5,rh=0.32)
    footer(s,2)

    # ===== 3. 전체 결과 =====
    s=slide(prs)
    header(s,"08 · 실험 결과","실험 결과 — 전체 결과",
           [("5개 모델 × 3 경로 연료 KPI (peak, 30회 평균). ",13.5,MUTED,False),("거리최단 < RL < 연료최단",13.5,INK,True),(" 경향.",13.5,MUTED,False)])
    rows=[["경로","① 거리최단","② 연료최단 TDD","③ rl_base","④ rl_signal","⑤ rl_attn"],
          [("OD-1","",True),"1715","1298 (−24.3%)","1435 (−16.3%)",("1393 (−18.8%)",ACCENT,True),"1502 (−12.4%)"],
          [("OD-2","",True),"1107","1014 (−8.4%)","1120 (+1.1%)","1154 (+4.3%)","1181 (+6.7%)"],
          [("OD-3","",True),"1368","1290 (−5.7%)","1357 (−0.8%)","1457 (+6.5%)","1468 (+7.3%)"],
          [("평균","",True),("1342",INK,True),("1194 (−11.1%)",GREEN,True),"1291 (−3.8%)","1322 (−1.5%)","1347 (+0.4%)"]]
    simple_table(s,0.4,1.85,12.5,[0.8,1.3,1.5,1.4,1.4,1.3],rows,fontsz=10.5,rh=0.52)
    rect(s,0.4,4.95,12.5,1.75,CARD,line=LINE,rounded=True)
    rect(s,0.55,5.1,0.06,0.3,ACCENT)
    text(s,0.7,5.06,11,0.3,[("실험 결과",11.5,ACCENT,True)])
    text(s,0.7,5.4,12.1,1.25,[
        [("연료최단 TDD가 거리최단 대비 평균 ",12.5,INK,False),("−11.1%",12.5,GREEN,True),(" 절감 — 강남 실데이터에서 '거리최소 ≠ 연료최소' 입증.",12.5,INK,False)],
        [("• OD-1(신논현)은 신호함정 뚜렷 → ",11.5,MUTED,False),("rl_signal −18.8%",11.5,INK,True),(" 로 오라클에 근접.",11.5,MUTED,False)],
        [("• OD-2·OD-3은 신호함정 약함 → RL은 도달하나 연료는 최단과 비슷(페널티로 '도달' 우선 학습).",11.5,MUTED,False)]],ls=1.25)
    footer(s,3)

    # ===== 4. 경로1 & 경로2 =====
    s=slide(prs)
    header(s,"08 · 실험 결과","실험 결과 — 경로1 & 경로2",
           [("신호 대기 많은 강남에서 ",13.5,MUTED,False),("신호 회피 경로",13.5,INK,True),("가 실제 연료 절감으로 이어지는가.",13.5,MUTED,False)])
    pic(s,FIG/"gn_od1_sinnonhyeon_suseo_fullmap.png",0.5,1.75,h=2.85)
    text(s,0.5,1.72,5.5,0.3,[("경로-1 신논현→수서",12.5,ACCENT,True)])
    simple_table(s,0.5,4.75,5.9,[1.5,1.0,1.0,0.8],[
        ["모델","연료","vs①","대기"],
        [("① 거리최단","",False),"1715","—","1174"],
        [("② 연료최단","",False),("1298",INK,True),("−24.3%",ACCENT,True),"285"],
        [("④ rl_signal","",False),"1393",("−18.8%",ACCENT,True),"229"],
        [("⑤ rl_attn","",False),"1502","−12.4%","278"]],fontsz=9.5,rh=0.34)
    pic(s,FIG/"gn_od2_yangjae_yeongdong_fullmap.png",7.0,1.75,h=2.85)
    text(s,7.0,1.72,5.5,0.3,[("경로-2 양재→영동대교",12.5,ACCENT,True)])
    simple_table(s,7.0,4.75,5.9,[1.5,1.0,1.0,0.8],[
        ["모델","연료","vs①","대기"],
        [("① 거리최단","",False),"1107","—","610"],
        [("② 연료최단","",False),("1014",INK,True),("−8.4%",ACCENT,True),"225"],
        [("④ rl_signal","",False),"1154","+4.3%","470"],
        [("⑤ rl_attn","",False),"1181","+6.7%","464"]],fontsz=9.5,rh=0.34)
    footer(s,4)

    # ===== 5. 경로3 (빈 자리 채움) =====
    s=slide(prs)
    header(s,"08 · 실험 결과","실험 결과 — 경로3",
           [("경로가 ",13.5,MUTED,False),("같아 보여도 다른 경로",13.5,INK,True),(" — 5.7% 연료차의 정체.",13.5,MUTED,False)])
    # 좌상: 결과해석
    rect(s,0.4,1.7,6.0,2.5,CARD,line=LINE,rounded=True)
    rect(s,0.55,1.85,0.06,0.3,ACCENT); text(s,0.7,1.81,5.6,0.3,[("결과 해석 — 경로-3 세곡→삼성",11.5,ACCENT,True)])
    text(s,0.7,2.2,5.55,1.95,[
        [("• 거리최단과 연료최단은 ",10.5,INK,False),("같은 corridor 공유",10.5,INK,True),("(공통 51노드)이나 ",10.5,INK,False),("중간 13노드가 다른 교차로",10.5,INK,True),(".",10.5,INK,False)],
        [("• 신호등 수 동일(12개), 우회 +92m(+1.2%) — 전체맵 줌엔 겹쳐 보임.",10.5,MUTED,False)],
        [("• 연료차 90mL 중 ",10.5,INK,False),("66mL(73%)가 공회전",10.5,ACC2,True),(" — fuel_TDD가 신호 녹색 타이밍 우회로 대기 631→507s(−124s).",10.5,INK,False)],
        [("→ '거리최소 ≠ 연료최소'가 같은 corridor 내 교차로 선택 미시 수준에서도 성립.",10.5,GREEN,True)]],ls=1.2)
    # 좌하: 표
    simple_table(s,0.4,4.4,6.0,[1.0,1.4,1.0,0.9,0.8],[
        ["OD-3","모델","연료","vs①","도달"],
        ["","① 거리최단","1368","—","100%"],
        ["",("② 연료최단",INK,True),("1290",INK,True),("−5.7%",ACCENT,True),"100%"],
        ["","③ rl_base","1357","−0.8%","93%"],
        ["","④ rl_signal","1457","+6.5%","100%"],
        ["","⑤ rl_attn","1468","+7.3%","100%"]],fontsz=9.5,rh=0.36)
    # 우: 5모델 경로 png
    pic(s,FIG/"gn_od3_allmodels_fullmap.png",6.7,1.7,h=5.0)
    text(s,6.7,6.72,6.3,0.25,[("OD-3 5개 모델 주행경로 — rl_signal(빨강)이 동측으로 분기",9,MUTED,False)])
    footer(s,5)

    # ===== 6. 일반화 검증 (NEW) =====
    s=slide(prs)
    header(s,"08 · 실험 결과","일반화 검증 — 처음 보는 12x12 (held-out)",
           [("다중 토폴로지 학습으로 ",13.5,MUTED,False),("전이 실패(0%)를 해결",13.5,INK,True),("했으나 ",13.5,MUTED,False),("강남 특화는 양보",13.5,ACC2,True),(".",13.5,MUTED,False)])
    # 좌: 12x12 held-out 표
    text(s,0.4,1.72,6.2,0.3,[("12x12 held-out (학습 미포함, 30회)",12.5,ACCENT,True)])
    simple_table(s,0.4,2.1,6.0,[1.6,1.4,1.0,1.0],[
        ["모델","연료(mL)","도달","vs단일학습"],
        [("① 거리최단","",False),"1226","100%","—"],
        [("② 연료최단","",False),("894",INK,True),"100%","—"],
        [("③ rl_base","",False),"3625","100%",("0%→100% ✓",GREEN,True)],
        [("④ rl_signal","",False),"2160","100%","100%→100%"],
        [("⑤ rl_attn","",False),("1408",ACCENT,True),"100%",("0%→100% ✓",GREEN,True)]],fontsz=9.5,rh=0.38)
    text(s,0.4,4.7,6.0,1.0,[
        [("• 단일학습 모델: base·attn ",10.5,INK,False),("12x12서 0% (제자리 회전)",10.5,ACC2,True)],
        [("• 다중학습 모델: 전부 ",10.5,INK,False),("100% 도달",10.5,GREEN,True),(" — 전이 성공.",10.5,INK,False)],
        [("• 신호밀집 12x12서 ",10.5,INK,False),("attn이 RL 최고(1408)",10.5,ACCENT,True),(" → 'attention 가설' 재현.",10.5,INK,False)]],ls=1.25)
    # 우: 트레이드오프
    rect(s,6.7,1.95,6.2,4.7,CARD,line=LINE,rounded=True)
    rect(s,6.85,2.1,0.06,0.3,ACC2); text(s,7.0,2.06,5.7,0.3,[("일반화 ↔ 특화 트레이드오프",12,ACC2,True)])
    simple_table(s,6.9,2.55,5.8,[2.0,1.5,1.5],[
        ["평가","단일학습","다중학습"],
        [("12x12 held-out","",True),("base/attn 0%",ACC2,False),("전부 100%",GREEN,True)],
        [("강남 3 OD","",True),("93~100%",GREEN,True),("0~10%",ACC2,True)]],fontsz=10,rh=0.5)
    text(s,6.9,4.5,5.8,2.0,[
        [("• 다중학습은 random-route로 강남 깊이 특화를 잃어 ",10.5,INK,False),("강남 긴 OD 퇴행",10.5,ACC2,True),(".",10.5,INK,False)],
        [("• 처음 보는 환경 강건성 ↑ ↔ in-domain 최적성 ↓.",10.5,MUTED,False)],
        [("→ 다음: 강남 fixed-OD 미세조정 병행, 또는 토폴로지별 head 분리로 둘 다 확보.",10.5,GREEN,True)]],ls=1.25)
    footer(s,6)

    # ===== 7. 기대효과 =====
    s=slide(prs)
    header(s,"08 · 실험 결과","기대효과",
           [("신호·연료 인식 경로탐색이 단순 최단거리 대비 연료 최대 ",13.5,MUTED,False),("20% 절감",13.5,ACC2,True),("함을 입증.",13.5,MUTED,False)])
    cards=[("01","신호 반영 경로탐색 = 16.6% 절감","정책",
            ["차량 2·5부제는 운행 강제 제한, 일 연료 3~4% 절감","본 모델은 경로만 바꿔 제한 없이 16.6% 절감"],"→ 정책 효과의 4~5배"),
           ("02","한화 약 1.3억 원/일 절감효과","경제",
            ["16.6%는 일 소비 59.3만 배럴 기준 약 10만 배럴","전체 차량 10%만 적용해도 1만 배럴 절감"],"→ 한화 약 1.3억 원/일"),
           ("03","운전 1회 1L, 연 10만 원 절감","개인",
            ["거리최단은 신호정체로 느림 → 우회로가 연료·시간 동시 단축","1회 1L = 연 약 10만 원(출퇴근 44회)"],"→ 더 빠르면서 기름값까지"),
           ("04","국내 도로망 전체 적용 — 전국 확장","기술",
            ["국가교통 데이터 사용 → 서울·전국 확장 가능","16.6% 절감률을 서울 통행량에 적용"],"→ 수송부문 에너지 절감 기여")]
    for i,(n,t,tag,bs,concl) in enumerate(cards):
        r,c=divmod(i,2); x=0.4+c*6.45; y=1.75+r*2.5
        rect(s,x,y,6.25,2.35,CARD,line=LINE,rounded=True)
        text(s,x+0.2,y+0.16,0.6,0.4,[(n,17,ACCENT,True)])
        text(s,x+0.85,y+0.2,4.6,0.4,[(t,12.5,INK,True)])
        text(s,x+5.3,y+0.2,0.85,0.3,[(tag,9.5,MUTED,True)],align=PP_ALIGN.RIGHT)
        yy=y+0.72
        for b in bs:
            text(s,x+0.25,yy,5.8,0.45,[("• ",10.5,ACCENT,True),(b,10.5,INK,False)],ls=1.0); yy+=0.46
        rect(s,x+0.25,y+1.92,5.75,0.32,ACCFILL,rounded=True,radius=0.3)
        text(s,x+0.4,y+1.92,5.6,0.32,[(concl,10.5,ACCENT,True)],anchor=MSO_ANCHOR.MIDDLE)
    footer(s,7)

    # ===== 8. 한계점 · 향후 과제 =====
    s=slide(prs)
    header(s,"08 · 실험 결과","한계점 · 향후 과제",
           [("본 연구의 사회·기술적 임팩트와 한계 · 향후 과제.",13.5,MUTED,False)])
    cards=[("1. 여러 환경 성능 보장?",[("속도: 링크 연결 구조 따라 상이한 결과 여지","",False),("신호: 밀도·연동 체계 지역마다 달라 일반화 난",ACC2,True),("→ 20% 같은 절대수치 일반 신뢰 어려움","",False)]),
           ("2. 실제 도로 변수 반영?",[("끼어들기·보행자·돌발 등 교통변수로 시나리오 주행 어려움","",False),("개별 차량 상호작용 미반영","",False),("→ 미시 교통변수 명시 시뮬 필요",GREEN,True)]),
           ("3. 경량 예측모델 VT-Micro",[("경량 모델로 가벼운 실행이 장점","",False),("실험 차량군 기반 → 모든 차종 보편 모델 아님","",False),("→ 경사·적재·날씨 반영 상위모델로 정확도↑",GREEN,True)])]
    for i,(t,bs) in enumerate(cards):
        x=0.4+i*4.25
        rect(s,x,1.75,4.05,3.7,CARD,line=LINE,rounded=True)
        text(s,x+0.2,1.9,3.7,0.34,[(t,12,ACCENT,True)])
        yy=2.45
        for (b,c,bold) in bs:
            if not isinstance(c,RGBColor): c=INK
            bullet=ACCENT if c is INK else c
            text(s,x+0.22,yy,3.65,0.8,[("• ",10.5,bullet,True),(b,10.5,c,bold)],ls=1.05); yy+=0.92
    rect(s,0.4,5.7,12.5,1.0,DARK,rounded=True)
    text(s,0.6,5.7,12.1,1.0,[("다양한 도로 환경 + 실제 교통 변수(차량 상호작용·보행자)를 반영한 추가 검증이 요구된다",15,ONDARK,True)],
         align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s,8)

    # ===== 9. 참고문헌 =====
    s=slide(prs)
    header(s,"REFERENCES","참고 문헌",None,None)
    refs=["[1] 에너지경제연구원 (2026.04.10). 「차량 운행 제한 및 재택근무 시행에 따른 석유제품 수요 절감 효과」. 보도자료.",
          "[2] 연합뉴스 (2026.03.11). 「기름값 싼 주유소에 차 수십 대 몰려…'원정 주유'도 성행」. 보도자료.",
          "[3] U.S. EPA (n.d.). Your Mileage May Vary. Green Vehicle Guide.",
          "[4] U.S. DOE (n.d.). Driving Habits. FuelEconomy.gov.",
          "[5] Wu et al. (2015). The Influence of Intersections on Fuel Consumption in Urban Road Traffic. PMC.",
          "[6] Davenport, C. (2021.05.18). Google Maps eco-friendly and safety features. XDA.",
          "[7] NACTO. Urban Street Design Guide — Intersection Design / Corner Radii.",
          "[8] FHWA. Handbook for Designing Roadways for the Aging Population, Jalali Khalilabadi et al. (2025)."]
    yy=1.95
    for r in refs:
        text(s,0.5,yy,12.3,0.4,[(r,11.5,INK,False)],ls=1.1); yy+=0.56
    footer(s,9)

    prs.save(str(OUT))
    print(f"saved → {OUT} ({len(prs.slides._sldIdLst)} slides)")

if __name__=="__main__":
    build()
