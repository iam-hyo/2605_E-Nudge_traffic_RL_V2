"""gen_ppt_resultpart.py
강남구 v2 실험 결과 PPT 생성 — PPT_style_guide.txt (ADDENDUM v3.1) 토큰 적용.

스타일 토큰:
  캔버스 16:9 13.333"×7.5" · 라이트 테마 · 폰트 Pretendard(폴백 Calibri)
  NAVY #1B2A4A / NAVY2 #24386B (사이드패널·강조박스)
  BLUE #2D6CDF (메인 액센트·번호뱃지) · RED #E23B3B (결론 뱃지 한정)
  연블루 #EAF1FC (칩) · 카드 #F5F8FC · 라인 #E2E7EF · INK #1F2937 · GRAY #5B6472
  본문 ≥14pt · 메타 12pt · 캡션 9~10pt

사용: venv/bin/python util/gen_ppt_resultpart.py <out_dir>
"""
from __future__ import annotations
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent

# ── 토큰 ──────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
NAVY2  = RGBColor(0x24, 0x38, 0x6B)
BLUE   = RGBColor(0x2D, 0x6C, 0xDF)
RED    = RGBColor(0xE2, 0x3B, 0x3B)
FILL   = RGBColor(0xEA, 0xF1, 0xFC)
CARD   = RGBColor(0xF5, 0xF8, 0xFC)
LINE   = RGBColor(0xE2, 0xE7, 0xEF)
INK    = RGBColor(0x1F, 0x29, 0x37)
GRAY   = RGBColor(0x5B, 0x64, 0x72)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x16, 0xA8, 0x58)

FONT = "Pretendard"   # 폴백: Calibri / Malgun Gothic (PowerPoint 자동 대체)
EMU_IN = 914400
CW, CH = 13.333, 7.5


def _set_font(run, size, color, bold=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 동아시아 폰트도 지정 (한글 글리프)
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': font})
    rPr.append(ea)


def _rect(slide, x, y, w, h, fill, line=None, line_w=1.0, shadow=False,
          rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=2, line_spacing=1.0):
    """runs: list of (text, size, color, bold) 또는 list of lists (멀티 문단)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if runs and not isinstance(runs[0], list):
        runs = [runs]
    for i, para_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para_runs:
            r = p.add_run(); r.text = txt
            _set_font(r, size, color, bold)
    return tb


def _chip(slide, x, y, w, label):
    c = _rect(slide, x, y, w, 0.34, FILL, line=BLUE, line_w=1.0, rounded=True)
    tf = c.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; _set_font(r, 10.5, BLUE, True)
    return c


def _slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _rect(s, 0, 0, CW, CH, WHITE)
    return s


def _header(slide, kicker, title):
    _rect(slide, 0, 0, 0.18, CH, BLUE)            # 좌측 액센트 바
    _text(slide, 0.55, 0.42, 12, 0.4, [(kicker, 12, BLUE, True)])
    _text(slide, 0.55, 0.74, 12.3, 0.7, [(title, 24, INK, True)])
    _rect(slide, 0.6, 1.46, 12.1, 0.018, LINE)


def _footer(slide, n):
    _text(slide, 0.55, 7.06, 8, 0.3,
          [("E-Nudge · 강남구 최소연료 경로탐색 RL · 2026-06-08", 9, GRAY, False)])
    _text(slide, 12.0, 7.06, 0.8, 0.3, [(f"{n:02d}", 9, GRAY, True)],
          align=PP_ALIGN.RIGHT)


# ── 슬라이드 빌더 ──────────────────────────────────────────────────────────────
def build(out_path: str, fig_dir: Path):
    prs = Presentation()
    prs.slide_width = Emu(int(CW * EMU_IN))
    prs.slide_height = Emu(int(CH * EMU_IN))

    # ===== 1. 표지 =====
    s = _slide(prs)
    _rect(s, 0, 0, 4.6, CH, NAVY)                  # 좌측 네이비 패널
    _rect(s, 0, 0, 4.6, 0.12, BLUE)
    _text(s, 0.5, 1.5, 3.7, 0.4, [("E-NUDGE PROJECT · RESULT PART", 12, RGBColor(0x9F,0xB3,0xD9), True)])
    _text(s, 0.5, 2.0, 3.7, 2.2, [
        [("강남구 실도로망", 27, WHITE, True)],
        [("최소연료 경로탐색", 27, WHITE, True)],
        [("강화학습 재실험 v2", 27, BLUE if False else RGBColor(0x8F,0xB4,0xF5), True)],
    ], line_spacing=1.05)
    _text(s, 0.5, 4.5, 3.7, 1.6, [
        [("신호·연료 인식 DQN으로", 13, RGBColor(0xC9,0xD4,0xE8), False)],
        [("\"가장 짧은 길이 가장 많은", 13, RGBColor(0xC9,0xD4,0xE8), False)],
        [(" 연료를 쓴다\"를 1995노드", 13, RGBColor(0xC9,0xD4,0xE8), False)],
        [(" 실데이터로 검증", 13, RGBColor(0xC9,0xD4,0xE8), False)],
    ], line_spacing=1.1)
    _text(s, 0.5, 6.7, 3.7, 0.4, [("2026-06-08 · output/08_1412_gangnam_v2", 10, RGBColor(0x7E,0x8E,0xB0), False)])

    # 우측 스탯 4칸 (2×2)
    stats = [
        ("−24.3%", "신논현→수서 peak 연료절감\n(연료최단 vs 거리최단)"),
        ("0% → 100%", "OD-3 RL 도달률\n(v1 실패 → v2 해결)"),
        ("900 runs", "3 OD × 2 시간대\n× 30회 × 5 모델"),
        ("1995 노드", "강남구 실도로망\n2439 링크 · 231 신호"),
    ]
    bx, by, bw, bh, gap = 5.05, 1.5, 3.75, 1.65, 0.3
    for i, (big, desc) in enumerate(stats):
        r, c = divmod(i, 2)
        x = bx + c * (bw + gap); y = by + r * (bh + gap)
        _rect(s, x, y, bw, bh, CARD, line=LINE, line_w=1.0, rounded=True)
        _text(s, x+0.25, y+0.2, bw-0.4, 0.6, [(big, 26, BLUE, True)])
        _text(s, x+0.25, y+0.82, bw-0.45, 0.75,
              [[(line, 11, GRAY, False)] for line in desc.split("\n")],
              line_spacing=1.05)
    # 역량 칩
    _text(s, 5.05, 5.15, 6, 0.3, [("CORE", 11, NAVY, True)])
    chips = ["환경 정비", "보상 설계", "동적 커리큘럼", "결과 분석", "시각화"]
    cx = 5.05
    for ch in chips:
        w = 0.42 + len(ch) * 0.165
        _chip(s, cx, 5.5, w, ch); cx += w + 0.18

    # ===== 2. 문제 & 환경 정비 =====
    s = _slide(prs)
    _header(s, "PROBLEM · ENVIRONMENT FIX", "진단된 4대 결함을 코드로 해소")
    _text(s, 0.6, 1.62, 12, 0.4,
          [("v1: Dijkstra 최단경로조차 통행 불가 → 에이전트 도달률 0%. 4개 결함을 정비해 학습 가능 환경으로 전환.", 13, GRAY, False)])
    rows = [
        ("①", "직진콘 ±5.7° 과소", "_movement_type 임계 0.1→0.5 (±30°)", "곡선도로 직진 인정"),
        ("②", "좌표축 뒤바뀜", "pos를 [경도(x), 위도(y)]로 정규화", "좌/우회전 부호 정상화"),
        ("③", "무신호 노드 좌회전 차단(93%)", "_node_allows_left 신호유무 최우선", "좌회전금지 93%→4.7%"),
        ("④", "Dijkstra 회전제한 미인지", "state=(node,prev)+turn-filter", "도달성 0%→99.3%"),
    ]
    y = 2.25
    for badge, cause, fix, eff in rows:
        _rect(s, 0.6, y, 12.1, 1.02, CARD, line=LINE, line_w=1.0, rounded=True)
        b = _rect(s, 0.8, y+0.26, 0.5, 0.5, BLUE, rounded=True)
        tf=b.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=badge; _set_font(r,18,WHITE,True)
        _text(s, 1.5, y+0.14, 4.2, 0.8, [
            [("원인", 9.5, RED, True)],
            [(cause, 13.5, INK, True)]], line_spacing=1.0)
        _text(s, 5.8, y+0.14, 4.0, 0.8, [
            [("대책", 9.5, BLUE, True)],
            [(fix, 12, INK, False)]], line_spacing=1.0)
        _text(s, 9.95, y+0.14, 2.6, 0.8, [
            [("효과", 9.5, GREEN, True)],
            [(eff, 12, GREEN, True)]], line_spacing=1.0)
        y += 1.14
    _footer(s, 2)

    # ===== 3. 실험 설계 — OD 3쌍 =====
    s = _slide(prs)
    _header(s, "EXPERIMENT DESIGN", "현실적 OD 3쌍 — 출근시간 통행 동선")
    od1_png = fig_dir / "gn_od1_sinnonhyeon_suseo_fullmap.png"
    if od1_png.exists():
        s.shapes.add_picture(str(od1_png), Inches(7.5), Inches(1.7),
                             height=Inches(5.0))
        _text(s, 7.5, 6.75, 5.5, 0.3,
              [("OD-1 신논현→수서 · 파랑=거리최단 주황=연료최단", 9.5, GRAY, False)])
    ods = [
        ("OD-1", "신논현 → 수서", "342695 → 613667",
         "강남대로 정체 통과 vs 우회 — 신호회피 효과 최대"),
        ("OD-2", "양재 → 영동대교", "338734 → 342812",
         "경부선 진입→영동대로 북상 한강도하 출근"),
        ("OD-3", "세곡 → 삼성(GBD)", "340301 → 341509",
         "주거지구→국제업무지구, 신호 밀집 도심 통과"),
    ]
    y = 1.85
    for tag, name, ids, why in ods:
        _rect(s, 0.6, y, 6.6, 1.42, CARD, line=LINE, line_w=1.0, rounded=True)
        _text(s, 0.85, y+0.16, 6.2, 0.4, [
            [(tag+"  ", 15, BLUE, True), (name, 15, INK, True)]])
        _text(s, 0.85, y+0.6, 6.2, 0.3, [(ids, 11, GRAY, False)])
        _text(s, 0.85, y+0.92, 6.2, 0.45, [(why, 11.5, INK, False)], line_spacing=1.0)
        y += 1.55
    _text(s, 0.6, 6.65, 6.6, 0.5,
          [("무작위 start/goal 폐지 · degree≥3·도달성 검증 · 동적 OD 커리큘럼 학습", 11, GRAY, False)],
          line_spacing=1.0)
    _footer(s, 3)

    # ===== 4. 핵심 결과 =====
    s = _slide(prs)
    _header(s, "KEY RESULTS", "OD-3 도달 해결 + 신논현 −24% 절감")
    # 좌: OD-3 도달 0→100
    _rect(s, 0.6, 1.75, 5.9, 2.4, CARD, line=LINE, line_w=1.0, rounded=True)
    _text(s, 0.85, 1.95, 5.4, 0.4, [("① OD-3 RL 도달률 — 완전 해결", 14, INK, True)])
    od3 = [("rl_base","0%","97% / 93%"),("rl_signal","0%","100% / 100%"),
           ("rl_attn","0%","100% / 100%")]
    yy=2.5
    for m,v1,v2 in od3:
        _text(s,0.95,yy,2.2,0.35,[(m,12,INK,False)])
        _text(s,3.1,yy,1.3,0.35,[(v1,12,RED,True)])
        _text(s,4.3,yy,0.5,0.35,[("→",12,GRAY,False)])
        _text(s,4.8,yy,1.6,0.35,[(v2,12,GREEN,True)])
        yy+=0.42
    _text(s,0.85,3.82,5.5,0.4,
          [("U턴 300 + 재방문 50 + 동적 부스트 → stub 루프 탈출", 10.5, GRAY, False)],line_spacing=1.0)
    # 우: 신논현 OD-1 절감
    _rect(s, 6.8, 1.75, 5.9, 2.4, CARD, line=LINE, line_w=1.0, rounded=True)
    _text(s, 7.05, 1.95, 5.4, 0.4, [("② 신논현→수서 peak 연료 (vs 거리최단)", 14, INK, True)])
    bars=[("② fuel_TDD",1298,"−24.3%",BLUE),("④ rl_signal",1393,"−18.8%",BLUE),
          ("③ rl_base",1435,"−16.3%",GRAY),("① shortest",1715,"—",RED)]
    yy=2.5; maxv=1715
    for m,val,pct,col in bars:
        _text(s,7.05,yy,2.1,0.3,[(m,11,INK,False)])
        bw=(val/maxv)*2.6
        _rect(s,9.2,yy+0.03,bw,0.22,col,rounded=False)
        _text(s,9.2+bw+0.1,yy-0.02,1.4,0.3,[(f"{val} ({pct})",10.5,col,True)])
        yy+=0.42
    # 하단 요약 바
    _rect(s, 0.6, 4.45, 12.1, 1.05, NAVY, rounded=True)
    _text(s, 0.95, 4.62, 11.6, 0.8, [
        [("종합 ", 13, WHITE, True),
         ("연료최단 TDD −11.1% · RL signal 도달케이스 −1.5% · 거리최단은 정체 직진으로 대기 1174s(최악)", 13, RGBColor(0xD6,0xE0,0xF2), False)],
        [("→ ", 13, RGBColor(0x8F,0xB4,0xF5), True),
         ("\"거리 최소화 ≠ 연료 최소화\"를 강남 1995노드 실데이터로 정량 입증", 13, WHITE, True)]],
        line_spacing=1.15)
    _footer(s, 4)

    # ===== 5. KPI 표 =====
    s = _slide(prs)
    _header(s, "KPI TABLE", "전 모델 × OD KPI (peak 08:00, 30회 평균)")
    hdr=["OD","모델","연료(mL)","vs 최단","대기(s)","도달"]
    data=[
        ["OD-1","① 거리최단","1715±49","—","1174","100%"],
        ["","② 연료최단 TDD","1298±38","−24.3%","285","100%"],
        ["","④ rl_signal","1393±85","−18.8%","229","100%"],
        ["","⑤ rl_attn","1502±155","−12.4%","278","97%"],
        ["OD-2","① 거리최단","1107±20","—","610","100%"],
        ["","② 연료최단 TDD","1014±39","−8.4%","225","100%"],
        ["","④ rl_signal","1154±135","+4.3%","470","100%"],
        ["OD-3","① 거리최단","1368±43","—","631","100%"],
        ["","② 연료최단 TDD","1290±46","−5.7%","507","100%"],
        ["","④ rl_signal","1457±204","+6.5%","606","100%"],
    ]
    rows=len(data)+1; cols=len(hdr)
    tbl_w=12.1; tbl_h=4.9
    gt=s.shapes.add_table(rows,cols,Inches(0.6),Inches(1.7),
                          Inches(tbl_w),Inches(tbl_h)).table
    widths=[1.3,3.2,2.0,1.9,1.8,1.9]
    for i,w in enumerate(widths): gt.columns[i].width=Inches(w*tbl_w/sum(widths))
    for j,h in enumerate(hdr):
        c=gt.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=NAVY
        c.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=h; _set_font(r,12,WHITE,True)
    for i,row in enumerate(data,1):
        shaded = "TDD" in row[1] or "signal" in row[1]
        for j,val in enumerate(row):
            c=gt.cell(i,j)
            c.fill.solid()
            c.fill.fore_color.rgb = FILL if (shaded and j>=2) else (CARD if i%2 else WHITE)
            c.vertical_anchor=MSO_ANCHOR.MIDDLE
            p=c.text_frame.paragraphs[0]
            p.alignment=PP_ALIGN.CENTER if j!=1 else PP_ALIGN.LEFT
            r=p.add_run(); r.text=val
            col = BLUE if (j==3 and val.startswith("−")) else (RED if (j==3 and val.startswith("+")) else INK)
            _set_font(r,11.5,col,bold=(j==3 and val not in("—","")))
    _footer(s, 5)

    # ===== 6. attention 역전 =====
    s = _slide(prs)
    _header(s, "ANALYSIS", "⭐ Attention 역전 — 12x12 최고 → 강남 최악")
    # 비교 표
    _rect(s,0.6,1.75,5.9,2.1,CARD,line=LINE,line_w=1.0,rounded=True)
    _text(s,0.85,1.92,5.4,0.35,[("RL 모델 순위 (vs 최단, 낮을수록 우수)",13,INK,True)])
    comp=[("",  "12x12","강남 v2"),
          ("1위","⑤attn −16.6%","③base −3.8%"),
          ("2위","④signal −13.3%","④signal −1.5%"),
          ("3위","③base −12.1%","⑤attn +0.4%")]
    yy=2.35
    for a,b,c in comp:
        bold=(a=="")
        _text(s,0.95,yy,1.0,0.3,[(a,11,GRAY,True)])
        _text(s,2.0,yy,2.3,0.3,[(b,11.5,BLUE if "attn" in b and a=="1위" else INK, "attn" in b)])
        _text(s,4.3,yy,2.1,0.3,[(c,11.5,RED if "attn" in c and a=="3위" else INK,"attn" in c)])
        yy+=0.36
    # 원인 4
    _text(s,6.8,1.78,6.0,0.4,[("역전 원인 4가지",14,INK,True)])
    causes=[
        ("① 과제 성격 변화(결정적)","12x12 신호 67%·신호함정=attn 홈그라운드 → 강남 신호 11.6%·토폴로지 주행성 문제"),
        ("② 예고된 리스크 현실화","12x12 §8: attn은 1.5× epoch 필요 → 동일 15000ep로 attn만 under-fit"),
        ("③ 페널티 지배 + 3-OD 분산","목표가 연료최소→루프회피·도달로 이동, 느린 attn 불리"),
        ("④ 강남은 극단함정 아님","fuel_TDD도 −11%(12x12 −20%), 미세최적화 여지 작음"),
    ]
    yy=2.25
    for t,d in causes:
        _text(s,6.8,yy,6.0,0.3,[(t,11.5,BLUE,True)])
        _text(s,6.8,yy+0.28,6.0,0.5,[(d,10.5,GRAY,False)],line_spacing=0.95)
        yy+=0.92
    # 한줄 요약 바
    _rect(s,0.6,6.05,12.1,0.92,NAVY,rounded=True)
    _text(s,0.95,6.18,11.6,0.7,[
        [("한 줄: ",13,RGBColor(0x8F,0xB4,0xF5),True),
         ("12x12는 attention의 시험문제(신호 식별), 강남은 robustness 시험문제(하드 그래프 항행).",13,WHITE,False)],
        [("        과제가 바뀌니 모델 순위가 뒤집혔다.",13,WHITE,True)]],line_spacing=1.1)
    _footer(s, 6)

    # ===== 7. 다음 방향성 =====
    s = _slide(prs)
    _header(s, "NEXT DIRECTIONS", "다음 방향성 — 모델 · 서빙 · 기대효과")
    cols3=[
        ("1  모델 개선", BLUE, [
            "보상 anneal (최우선) — 페널티 점진감쇠로 연료최적성 회복",
            "attention 공정 재학습 (1.5~2× epoch)",
            "sparse-attention — 신호 노드 한정",
            "fuel-potential shaping",
            "단계 커리큘럼 정식화",
        ]),
        ("2  실제 서빙", NAVY2, [
            "RL 1-step 추론(ms) → 실시간 라우팅 (TDD는 전체 Dijkstra라 느림)",
            "모델-프리 → 실데이터 미비 지역 강점",
            "선결: 일방통행·회전제한·SPAT 정합",
            "임의 OD 일반화 검증",
            "TMAP/Kakao A/B → 에코모드",
        ]),
        ("3  기대효과·한계", RED, [
            "[+] 출근첨두 연료 −24%(오라클)·−19%(RL)",
            "[+] 실 도로망 RL 학습 성립 증명",
            "[−] 신호회피 이득 첨두·OD 의존",
            "[−] 페널티-연료최적성 트레이드오프",
            "[−] 실도로 데이터 정합 미완",
        ]),
    ]
    bw=3.9; bx=0.6; gap=0.2
    for i,(title,col,items) in enumerate(cols3):
        x=bx+i*(bw+gap)
        _rect(s,x,1.7,bw,5.0,CARD,line=LINE,line_w=1.0,rounded=True)
        _rect(s,x,1.7,bw,0.62,col,rounded=False)
        tf=_text(s,x+0.2,1.78,bw-0.4,0.5,[(title,15,WHITE,True)])
        yy=2.55
        for it in items:
            _text(s,x+0.25,yy,bw-0.45,0.7,[
                [("• ",11,col,True),(it,11,INK,False)]],line_spacing=0.98)
            yy+=0.84
    _footer(s, 7)

    # ===== 8. 결론 =====
    s = _slide(prs)
    _rect(s,0,0,CW,CH,WHITE)
    _rect(s,0,0,CW,1.0,NAVY)
    _text(s,0.6,0.28,12,0.5,[("CONCLUSION",13,RGBColor(0x8F,0xB4,0xF5),True),
                             ("   결론",20,WHITE,True)])
    _text(s,0.7,1.5,11.9,1.2,[
        [("환경 정비 4종 + 보상 설계 + 동적 OD 커리큘럼으로,",18,INK,True)],
        [("v1의 도달률 0% 실패를 ",18,INK,False),("도달 93~100%의 비교가능 실험",18,BLUE,True),
         ("으로 전환했다.",18,INK,False)]],line_spacing=1.2)
    pts=[
        ("신호·연료 인식 경로탐색 가설","강남 1995노드 실데이터에서 성립 — 신논현 OD-1 연료 −24%(오라클)·−19%(RL)"),
        ("도달성-연료최적성 트레이드오프","페널티로 OD-3 도달 확보, 연료최적성 일부 양보 → 보상 anneal이 다음 과제"),
        ("Attention 규모 확장성 한계","신호 sparse 대규모에서 단순 모델 대비 열위 → 과제 성격이 모델 우열을 결정"),
    ]
    yy=3.1
    for i,(t,d) in enumerate(pts,1):
        b=_rect(s,0.7,yy,0.5,0.5,RED if i==3 else BLUE,rounded=True)
        tf=b.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=str(i); _set_font(r,16,WHITE,True)
        _text(s,1.4,yy+0.02,11.3,0.7,[
            [(t,14.5,INK,True)],
            [(d,12,GRAY,False)]],line_spacing=1.0)
        yy+=1.05
    _rect(s,0,7.1,CW,0.4,NAVY)
    _text(s,0.6,7.16,12,0.3,[("E-Nudge Project · 강남구 RL 재실험 v2 · output/08_1412_gangnam_v2",10,WHITE,False)])

    prs.save(out_path)
    print(f"  saved → {out_path}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    out_dir = Path(sys.argv[1]) if len(sys.argv) > 1 else (ROOT/"output"/"08_1412_gangnam_v2")
    out_dir.mkdir(parents=True, exist_ok=True)
    build(str(out_dir/"E-nudge_ResultPart_gangnam_v2.pptx"), out_dir)
